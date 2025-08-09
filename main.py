import os
import io
import argparse
from dotenv import load_dotenv
import google.generativeai as genai
from pptx import Presentation
from PIL import Image
import pytesseract

# --- Configuration ---
# Load environment variables from .env file
load_dotenv()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

def configure_ai():
    """Configures the Gemini API."""
    if not GEMINI_API_KEY:
        raise ValueError("GEMINI_API_KEY not found. Please set it in your .env file.")
    genai.configure(api_key=GEMINI_API_KEY)

def extract_text_from_image(image_bytes):
    """
    Performs OCR on an image to extract text.
    """
    try:
        image = Image.open(io.BytesIO(image_bytes))
        text = pytesseract.image_to_string(image)
        return text.strip()
    except Exception as e:
        print(f"    - Warning: Could not process an image. Error: {e}")
        return ""

def extract_data_from_pptx(pptx_path):
    """
    Extracts all textual content from a PowerPoint presentation.
    """
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"The file '{pptx_path}' was not found.")

    prs = Presentation(pptx_path)
    slide_data = {}
    print(f"[+] Extracting text from {len(prs.slides)} slides...")

    for i, slide in enumerate(prs.slides):
        slide_number = i + 1
        slide_texts = []
        image_count = 0

        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_texts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        slide_texts.append(cell.text_frame.text)
            
            if hasattr(shape, 'image'):
                image_count += 1
                image_bytes = shape.image.blob
                ocr_text = extract_text_from_image(image_bytes)
                if ocr_text:
                    slide_texts.append(f"[OCR Text from Image]:\n{ocr_text}")
        
        slide_data[slide_number] = "\n".join(filter(None, slide_texts))
        
        if image_count > 0:
            print(f"    - Extracted text and {image_count} image(s) from Slide {slide_number}")
        else:
            print(f"    - Extracted text from Slide {slide_number}")

    return slide_data

def analyze_text_with_gemini(all_slides_text):
    """
    Sends the consolidated presentation text to the Gemini API for analysis.
    """
    model = genai.GenerativeModel('gemini-1.5-flash-latest')
    
    prompt = f"""
    You are an expert business and strategy consultant for a top-tier firm like McKinsey or BCG. 
    Your task is to meticulously analyze the following content extracted from a multi-slide PowerPoint presentation. 
    The content from each slide is clearly separated and labeled with '--- Slide X ---'.

    Your goal is to identify all factual, numerical, and logical inconsistencies across the entire deck. Pay close attention to:
    1.  **Conflicting Numerical Data:** Look for mismatched revenue figures, user counts, market shares, financial projections, or percentages that don't add up. Be precise. For example, '$15M' on one slide vs. '$12M' on another.
    2.  **Contradictory Textual Claims:** Identify statements that contradict each other. For example, "we are entering a low-competition market" on one slide and "we must navigate a highly competitive landscape" on another.
    3.  **Timeline Mismatches:** Find conflicting dates, project phases, or launch forecasts.

    **Instructions for Output:**
    - Analyze the content below thoroughly.
    - If you find one or more inconsistencies, list each one clearly. For each inconsistency, provide:
        - A short, bolded title for the inconsistency (e.g., **Conflicting Revenue Projections**).
        - The slide numbers where the conflicting information appears (e.g., "Slide 2 vs. Slide 4.").
        - The specific conflicting data or statements, quoting them if possible.
        - A brief explanation of why it's an inconsistency.
    - If you find absolutely no inconsistencies, you MUST respond with the single sentence: "No inconsistencies were found in the presentation."
    - Do not comment on the presentation's quality, style, or grammar. Focus ONLY on factual and logical contradictions.

    Here is the presentation content:
    ---
    {all_slides_text}
    ---
    """
    
    try:
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"An error occurred while communicating with the AI model: {e}"

def main():
    """Main function to run the DeckScanner tool."""
    parser = argparse.ArgumentParser(description="Analyze a PowerPoint presentation for inconsistencies using AI.")
    parser.add_argument("pptx_file", help="Path to the .pptx file to analyze.")
    args = parser.parse_args()

    try:
        configure_ai()
        print(f"[+] Starting analysis of '{args.pptx_file}'...")
        
        slide_data = extract_data_from_pptx(args.pptx_file)
        
        print("[+] Extraction complete. Consolidating text for analysis.")
        consolidated_text = "\n\n".join(
            f"--- Slide {num} ---\n{text}" for num, text in slide_data.items()
        )
        
        print("[+] Sending data to AI for inconsistency analysis...")
        ai_report = analyze_text_with_gemini(consolidated_text)
        
        print("[+] Analysis Complete.\n")
        print("========================================")
        print("      AI Inconsistency Report")
        print("========================================")
        print(ai_report)

    except (FileNotFoundError, ValueError) as e:
        print(f"[!] Error: {e}")
    except Exception as e:
        print(f"[!] An unexpected error occurred: {e}")

if __name__ == "__main__":
    main()